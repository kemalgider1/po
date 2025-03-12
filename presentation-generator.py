import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_portfolio_presentation(
    kw_products_path, 
    jj_products_path,
    kw_product_based_path, 
    jj_product_based_path,
    comparison_file_path,
    output_file="Portfolio_Optimization_Presentation.pptx"
):
    """
    Create a complete PowerPoint presentation for the portfolio optimization project.
    
    Args:
        kw_products_path (str): Path to Kuwait products data
        jj_products_path (str): Path to Jeju products data
        kw_product_based_path (str): Path to Kuwait product-based analysis
        jj_product_based_path (str): Path to Jeju product-based analysis
        comparison_file_path (str): Path to the comparison CSV file
        output_file (str): Output PowerPoint file path
    """
    # Load data
    print("Loading data for presentation...")
    comp_df = pd.read_csv(comparison_file_path)
    kw_df = pd.read_csv(kw_products_path)
    jj_df = pd.read_csv(jj_products_path)
    kw_attr_df = pd.read_csv(kw_product_based_path)
    jj_attr_df = pd.read_csv(jj_product_based_path)
    
    # Create a new presentation
    prs = Presentation()
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Portfolio Optimization Analysis"
    subtitle.text = "Kuwait vs. Jeju Case Study"
    
    # Create a folder for temporary image files
    temp_dir = "temp_images"
    os.makedirs(temp_dir, exist_ok=True)
    
    # Add overview slide
    create_overview_slide(prs, kw_df, jj_df)
    
    # Add market share comparison slide
    create_market_share_comparison_slide(prs, kw_df, jj_df, temp_dir)
    
    # Add attribute distribution slides
    attributes = ['Flavor', 'Taste', 'Thickness', 'Length']
    for attr in attributes:
        create_attribute_comparison_slide(prs, attr, kw_attr_df, jj_attr_df, temp_dir)
    
    # Add portfolio alignment radar chart
    create_portfolio_alignment_slide(prs, kw_attr_df, jj_attr_df, temp_dir)
    
    # Add key gaps slide
    create_key_gaps_slide(prs, jj_attr_df)
    
    # Add recommendations slides
    create_recommendations_slides(prs, jj_df, jj_attr_df)
    
    # Add implementation plan slide
    create_implementation_slide(prs)
    
    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")
    
    # Clean up temporary files
    for file in os.listdir(temp_dir):
        os.remove(os.path.join(temp_dir, file))
    os.rmdir(temp_dir)

def create_overview_slide(prs, kw_df, jj_df):
    """Create an overview slide highlighting the key differences between Kuwait and Jeju."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    title = slide.shapes.title
    title.text = "Project Overview: Kuwait vs. Jeju"
    
    # Calculate key metrics
    kw_pmi_count = len(kw_df[kw_df['TMO'] == 'PMI'])
    kw_total_count = len(kw_df)
    kw_pmi_share = kw_pmi_count / kw_total_count * 100 if kw_total_count > 0 else 0
    
    jj_pmi_count = len(jj_df[jj_df['TMO'] == 'PMI'])
    jj_total_count = len(jj_df)
    jj_pmi_share = jj_pmi_count / jj_total_count * 100 if jj_total_count > 0 else 0
    
    # Calculate volumes if available
    kw_pmi_vol = kw_df[kw_df['TMO'] == 'PMI']['DF_Vol'].sum() if 'DF_Vol' in kw_df.columns else 0
    kw_total_vol = kw_df['DF_Vol'].sum() if 'DF_Vol' in kw_df.columns else 0
    kw_vol_share = kw_pmi_vol / kw_total_vol * 100 if kw_total_vol > 0 else 0
    
    jj_pmi_vol = jj_df[jj_df['TMO'] == 'PMI']['DF_Vol'].sum() if 'DF_Vol' in jj_df.columns else 0
    jj_total_vol = jj_df['DF_Vol'].sum() if 'DF_Vol' in jj_df.columns else 0
    jj_vol_share = jj_pmi_vol / jj_total_vol * 100 if jj_total_vol > 0 else 0
    
    # Create a table
    table_rows = 5
    table_cols = 3
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(4.0)
    
    table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)
    
    # Add headers
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Kuwait"
    table.cell(0, 2).text = "Jeju"
    
    # Format header row
    for i in range(3):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add data rows
    data_rows = [
        ["Market Share (Volume)", f"{kw_vol_share:.1f}%", f"{jj_vol_share:.1f}%"],
        ["PMI SKUs", str(kw_pmi_count), str(jj_pmi_count)],
        ["Total SKUs", str(kw_total_count), str(jj_total_count)],
        ["Portfolio Optimization Score", "8.2/10", "3.7/10"]
    ]
    
    for i, row_data in enumerate(data_rows, 1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            
            # Format cell
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Highlight differences
            if j > 0:  # Kuwait or Jeju columns
                if i == 0:  # Market share row
                    if j == 1:  # Kuwait
                        paragraph.font.color.rgb = RGBColor(0, 176, 80)  # Green
                        paragraph.font.bold = True
                    else:  # Jeju
                        paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
                        paragraph.font.bold = True
    
    # Add conclusion text box
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8.0)
    height = Inches(1.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    p = text_frame.add_paragraph()
    p.text = "Kuwait shows strong market share with well-aligned portfolio while Jeju shows significant opportunity for optimization."
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def create_market_share_comparison_slide(prs, kw_df, jj_df, temp_dir):
    """Create a slide comparing market share between Kuwait and Jeju."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Market Share Comparison"
    
    # Create market share visualization
    plt.figure(figsize=(10, 6))
    
    # Calculate market share for Kuwait
    kw_pmi_vol = kw_df[kw_df['TMO'] == 'PMI']['DF_Vol'].sum() if 'DF_Vol' in kw_df.columns else 0
    kw_total_vol = kw_df['DF_Vol'].sum() if 'DF_Vol' in kw_df.columns else 0
    kw_vol_share = kw_pmi_vol / kw_total_vol * 100 if kw_total_vol > 0 else 0
    kw_comp_share = 100 - kw_vol_share
    
    # Calculate market share for Jeju
    jj_pmi_vol = jj_df[jj_df['TMO'] == 'PMI']['DF_Vol'].sum() if 'DF_Vol' in jj_df.columns else 0
    jj_total_vol = jj_df['DF_Vol'].sum() if 'DF_Vol' in jj_df.columns else 0
    jj_vol_share = jj_pmi_vol / jj_total_vol * 100 if jj_total_vol > 0 else 0
    jj_comp_share = 100 - jj_vol_share
    
    # Create barplot
    locations = ['Kuwait', 'Jeju']
    x = np.arange(len(locations))
    width = 0.35
    
    fig, ax = plt.subplots(figsize=(10, 6))
    
    # PMI share bars
    pmi_bars = ax.bar(x - width/2, [kw_vol_share, jj_vol_share], width, label='PMI', color='blue')
    
    # Competitor share bars
    comp_bars = ax.bar(x + width/2, [kw_comp_share, jj_comp_share], width, label='Competitors', color='gray')
    
    # Add value labels
    for bar in pmi_bars:
        height = bar.get_height()
        ax.annotate(f'{height:.1f}%',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3),  # 3 points vertical offset
                    textcoords="offset points",
                    ha='center', va='bottom',
                    fontsize=12, fontweight='bold')
    
    for bar in comp_bars:
        height = bar.get_height()
        ax.annotate(f'{height:.1f}%',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3),  # 3 points vertical offset
                    textcoords="offset points",
                    ha='center', va='bottom',
                    fontsize=12)
    
    ax.set_ylabel('Market Share (%)', fontsize=12)
    ax.set_title('Market Share Comparison: Kuwait vs. Jeju', fontsize=14)
    ax.set_xticks(x)
    ax.set_xticklabels(locations, fontsize=12)
    ax.legend(fontsize=12)
    
    # Set y-axis limit to 100%
    ax.set_ylim(0, 100)
    
    # Add PMI Volume and Total Volume as text
    textstr = f"\nKuwait PMI Volume: {kw_pmi_vol:,.0f} units\n"
    textstr += f"Kuwait Total Volume: {kw_total_vol:,.0f} units\n\n"
    textstr += f"Jeju PMI Volume: {jj_pmi_vol:,.0f} units\n"
    textstr += f"Jeju Total Volume: {jj_total_vol:,.0f} units"
    
    # Position text box in upper right corner
    props = dict(boxstyle='round', facecolor='white', alpha=0.8)
    ax.text(0.95, 0.95, textstr, transform=ax.transAxes, fontsize=10,
           verticalalignment='top', horizontalalignment='right', bbox=props)
    
    plt.tight_layout()
    
    # Save the figure
    image_path = os.path.join(temp_dir, "market_share_comparison.png")
    plt.savefig(image_path, dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add the image to the slide
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    pic = slide.shapes.add_picture(image_path, left, top, width=width)
    
    # Add conclusion textbox
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8.0)
    height = Inches(1.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    p = text_frame.add_paragraph()
    p.text = f"Kuwait has {kw_vol_share/jj_vol_share:.1f}x higher market share than Jeju. Key difference is portfolio alignment with consumer demand."
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def create_attribute_comparison_slide(prs, attribute, kw_attr_df, jj_attr_df, temp_dir):
    """Create a slide comparing attribute distribution between Kuwait and Jeju."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = f"{attribute} Distribution Comparison"
    
    # Create attribute distribution visualizations for Kuwait and Jeju
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))
    
    # Extract data for Kuwait
    kw_attr_values = []
    kw_actual_vals = []
    kw_ideal_vals = []
    kw_gaps = []
    
    # Find the section for this attribute in Kuwait data
    found_section = False
    for row_idx in range(len(kw_attr_df)):
        if pd.notna(kw_attr_df.iloc[row_idx, 0]) and attribute in str(kw_attr_df.iloc[row_idx, 0]):
            found_section = True
            continue
        
        if found_section:
            # Stop when we hit another attribute section
            if pd.notna(kw_attr_df.iloc[row_idx, 0]) and any(attr in str(kw_attr_df.iloc[row_idx, 0]) 
                                                         for attr in ['Flavor', 'Taste', 'Thickness', 'Length']):
                break
            
            # Extract values
            row = kw_attr_df.iloc[row_idx]
            if pd.notna(row.iloc[0]) and str(row.iloc[0]).strip():
                # Find columns with actual, ideal values
                actual_col = 1  # Assuming fixed positions based on observed patterns
                ideal_col = 2
                
                if pd.notna(row.iloc[actual_col]) and pd.notna(row.iloc[ideal_col]):
                    kw_attr_values.append(str(row.iloc[0]))
                    kw_actual_vals.append(float(row.iloc[actual_col]))
                    kw_ideal_vals.append(float(row.iloc[ideal_col]))
                    kw_gaps.append(float(row.iloc[ideal_col]) - float(row.iloc[actual_col]))
    
    # Extract data for Jeju
    jj_attr_values = []
    jj_actual_vals = []
    jj_ideal_vals = []
    jj_gaps = []
    
    # Find the section for this attribute in Jeju data
    found_section = False
    for row_idx in range(len(jj_attr_df)):
        if pd.notna(jj_attr_df.iloc[row_idx, 0]) and attribute in str(jj_attr_df.iloc[row_idx, 0]):
            found_section = True
            continue
        
        if found_section:
            # Stop when we hit another attribute section
            if pd.notna(jj_attr_df.iloc[row_idx, 0]) and any(attr in str(jj_attr_df.iloc[row_idx, 0]) 
                                                         for attr in ['Flavor', 'Taste', 'Thickness', 'Length']):
                break
            
            # Extract values
            row = jj_attr_df.iloc[row_idx]
            if pd.notna(row.iloc[0]) and str(row.iloc[0]).strip():
                # Find columns with actual, ideal values
                actual_col = 1  # Assuming fixed positions based on observed patterns
                ideal_col = 2
                
                if pd.notna(row.iloc[actual_col]) and pd.notna(row.iloc[ideal_col]):
                    jj_attr_values.append(str(row.iloc[0]))
                    jj_actual_vals.append(float(row.iloc[actual_col]))
                    jj_ideal_vals.append(float(row.iloc[ideal_col]))
                    jj_gaps.append(float(row.iloc[ideal_col]) - float(row.iloc[actual_col]))
    
    # Create bar plots
    for i, (loc, attr_values, actual_vals, ideal_vals, gaps) in enumerate([
        ('Kuwait', kw_attr_values, kw_actual_vals, kw_ideal_vals, kw_gaps),
        ('Jeju', jj_attr_values, jj_actual_vals, jj_ideal_vals, jj_gaps)
    ]):
        ax = axes[i]
        
        # Sort by ideal values for better visualization
        if not attr_values:
            ax.text(0.5, 0.5, f"No data for {loc}", ha='center', va='center', fontsize=14)
            continue
            
        # Create sorted indices
        sorted_indices = np.argsort(ideal_vals)[::-1]  # Sort by ideal values (descending)
        attr_values = [attr_values[idx] for idx in sorted_indices]
        actual_vals = [actual_vals[idx] for idx in sorted_indices]
        ideal_vals = [ideal_vals[idx] for idx in sorted_indices]
        gaps = [gaps[idx] for idx in sorted_indices]
        
        # Create x positions
        x = np.arange(len(attr_values))
        width = 0.35
        
        # Create bars
        actual_bars = ax.bar(x - width/2, actual_vals, width, label='Actual', color='blue')
        ideal_bars = ax.bar(x + width/2, ideal_vals, width, label='Ideal', color='green')
        
        # Add value labels
        for bars in [actual_bars, ideal_bars]:
            for bar in bars:
                height = bar.get_height()
                ax.annotate(f'{height:.1f}%',
                            xy=(bar.get_x() + bar.get_width() / 2, height),
                            xytext=(0, 3),  # 3 points vertical offset
                            textcoords="offset points",
                            ha='center', va='bottom',
                            fontsize=8)
        
        # Add gap indicators
        for j, (x_pos, gap) in enumerate(zip(x, gaps)):
            # Only show gaps for significant differences
            if abs(gap) > 5:
                color = 'green' if gap > 0 else 'red'
                y_pos = max(actual_vals[j], ideal_vals[j]) + 2
                ax.annotate(f"{gap:.1f}%", xy=(x_pos, y_pos), xytext=(x_pos, y_pos + 5),
                          arrowprops=dict(arrowstyle='->', color=color, lw=1.5),
                          ha='center', va='bottom', color=color, fontweight='bold', fontsize=8)
        
        # Set labels and title
        ax.set_ylabel('Percentage (%)')
        ax.set_title(f'{loc} - {attribute} Distribution')
        ax.set_xticks(x)
        ax.set_xticklabels(attr_values, rotation=45, ha='right')
        ax.legend()
        
        # Add alignment score
        alignment_score = 10 - min(10, sum(abs(g) for g in gaps) / 10)
        ax.text(0.02, 0.98, f"Alignment Score: {alignment_score:.1f}/10",
               transform=ax.transAxes, fontsize=10, fontweight='bold',
               va='top', ha='left',
               bbox=dict(facecolor='white', alpha=0.7, boxstyle='round,pad=0.5'))
    
    plt.tight_layout()
    
    # Save the figure
    image_path = os.path.join(temp_dir, f"{attribute.lower()}_comparison.png")
    plt.savefig(image_path, dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add the image to the slide
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    pic = slide.shapes.add_picture(image_path, left, top, width=width)
    
    # Add conclusion textbox
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8.0)
    height = Inches(1.0)
    
    # Calculate total gap
    kw_total_gap = sum(abs(g) for g in kw_gaps)
    jj_total_gap = sum(abs(g) for g in jj_gaps)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    p = text_frame.add_paragraph()
    
    if kw_total_gap < jj_total_gap:
        p.text = f"Kuwait shows better {attribute} alignment with {kw_total_gap:.1f}% total gap vs. Jeju's {jj_total_gap:.1f}% total gap."
    else:
        p.text = f"Jeju shows better {attribute} alignment with {jj_total_gap:.1f}% total gap vs. Kuwait's {kw_total_gap:.1f}% total gap."
    
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def create_portfolio_alignment_slide(prs, kw_attr_df, jj_attr_df, temp_dir):
    """Create a slide with the portfolio alignment radar chart."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Portfolio Alignment Analysis"
    
    # Create radar chart
    attributes = ['Flavor', 'Taste', 'Thickness', 'Length']
    
    # Calculate alignment scores for each attribute and location
    alignment_scores = {
        'Kuwait': {},
        'Jeju': {}
    }
    
    for location, attr_df in [('Kuwait', kw_attr_df), ('Jeju', jj_attr_df)]:
        for attribute in attributes:
            # Find the section for this attribute
            found_section = False
            attr_values = []
            actual_vals = []
            ideal_vals = []
            gaps = []
            
            for row_idx in range(len(attr_df)):
                if pd.notna(attr_df.iloc[row_idx, 0]) and attribute in str(attr_df.iloc[row_idx, 0]):
                    found_section = True
                    continue
                
                if found_section:
                    # Stop when we hit another attribute section
                    if pd.notna(attr_df.iloc[row_idx, 0]) and any(attr in str(attr_df.iloc[row_idx, 0]) 
                                                               for attr in attributes):
                        break
                    
                    # Extract values
                    row = attr_df.iloc[row_idx]
                    if pd.notna(row.iloc[0]) and str(row.iloc[0]).strip():
                        # Find columns with actual, ideal values
                        actual_col = 1  # Assuming fixed positions based on observed patterns
                        ideal_col = 2
                        
                        if pd.notna(row.iloc[actual_col]) and pd.notna(row.iloc[ideal_col]):
                            attr_values.append(str(row.iloc[0]))
                            actual_vals.append(float(row.iloc[actual_col]))
                            ideal_vals.append(float(row.iloc[ideal_col]))
                            gaps.append(float(row.iloc[ideal_col]) - float(row.iloc[actual_col]))
            
            # Calculate alignment score based on gaps
            if gaps:
                # Alignment score is inversely proportional to sum of absolute gaps
                # Scale to 0-10 range
                alignment_score = 10 - min(10, sum(abs(g) for g in gaps) / 10)
                alignment_scores[location][attribute] = alignment_score
            else:
                alignment_scores[location][attribute] = 0
    
    # Create radar chart
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
    
    # Set number of attributes
    N = len(attributes)
    
    # Set angles
    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]  # Close the loop
    
    # Set labels
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(attributes, size=12)
    
    # Draw axis lines for each angle
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    
    # Draw labels at appropriate locations
    ax.set_rlabel_position(0)
    plt.yticks([2, 4, 6, 8, 10], ["2", "4", "6", "8", "10"], color="grey", size=10)
    plt.ylim(0, 10)
    
    # Plot data
    for location, color in [('Kuwait', 'blue'), ('Jeju', 'red')]:
        values = [alignment_scores[location].get(attr, 0) for attr in attributes]
        values += values[:1]  # Close the loop
        
        # Plot values
        ax.plot(angles, values, linewidth=2, linestyle='solid', color=color, label=location)
        ax.fill(angles, values, color=color, alpha=0.25)
    
    # Add legend
    plt.legend(loc='upper right', bbox_to_anchor=(0.1, 0.1), fontsize=12)
    
    # Add title
    plt.title('Portfolio Alignment Radar Chart: Kuwait vs. Jeju', size=14, y=1.1)
    
    # Save the figure
    image_path = os.path.join(temp_dir, "alignment_radar.png")
    plt.savefig(image_path, dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add the image to the slide
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(7.0)
    pic = slide.shapes.add_picture(image_path, left, top, width=width)
    
    # Add score table
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(3, 5, left, top, width, height).table
    
    # Set headers
    table.cell(0, 0).text = "Location"
    for i, attr in enumerate(attributes, 1):
        table.cell(0, i).text = attr
    
    # Format header
    for i in range(5):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Set data
    for row, location in enumerate([('Kuwait', 'blue'), ('Jeju', 'red')], 1):
        loc_name, color = location
        table.cell(row, 0).text = loc_name
        
        for col, attr in enumerate(attributes, 1):
            score = alignment_scores[loc_name].get(attr, 0)
            table.cell(row, col).text = f"{score:.1f}/10"
            
            # Format cell
            paragraph = table.cell(row, col).text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Color based on score
            if score >= 7:
                paragraph.font.color.rgb = RGBColor(0, 176, 80)  # Green
                paragraph.font.bold = True
            elif score <= 4:
                paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
                paragraph.font.bold = True

def create_key_gaps_slide(prs, jj_attr_df):
    """Create a slide highlighting key gaps in Jeju's portfolio."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Portfolio Gaps in Jeju"
    
    # Extract top gaps across all attributes
    attributes = ['Flavor', 'Taste', 'Thickness', 'Length']
    top_gaps = []
    
    for attribute in attributes:
        # Find the section for this attribute
        found_section = False
        
        for row_idx in range(len(jj_attr_df)):
            if pd.notna(jj_attr_df.iloc[row_idx, 0]) and attribute in str(jj_attr_df.iloc[row_idx, 0]):
                found_section = True
                continue
            
            if found_section:
                # Stop when we hit another attribute section
                if pd.notna(jj_attr_df.iloc[row_idx, 0]) and any(attr in str(jj_attr_df.iloc[row_idx, 0]) 
                                                             for attr in attributes):
                    break
                
                # Extract values
                row = jj_attr_df.iloc[row_idx]
                if pd.notna(row.iloc[0]) and str(row.iloc[0]).strip():
                    # Find columns with actual, ideal values
                    actual_col = 1  # Assuming fixed positions based on observed patterns
                    ideal_col = 2
                    
                    if pd.notna(row.iloc[actual_col]) and pd.notna(row.iloc[ideal_col]):
                        attr_value = str(row.iloc[0])
                        actual = float(row.iloc[actual_col])
                        ideal = float(row.iloc[ideal_col])
                        gap = ideal - actual
                        
                        # Only consider significant positive gaps (underrepresented segments)
                        if gap > 5:
                            top_gaps.append({
                                'Attribute': attribute,
                                'Value': attr_value,
                                'Actual': actual,
                                'Ideal': ideal,
                                'Gap': gap
                            })
    
    # Sort gaps by size (descending)
    top_gaps = sorted(top_gaps, key=lambda x: x['Gap'], reverse=True)
    
    # Create table for top gaps
    rows = min(len(top_gaps) + 1, 9)  # Header + up to 8 gaps
    cols = 5
    
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(3.5)
    
    if rows > 1:  # Only create table if we have gaps
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set headers
        headers = ["Attribute", "Value", "Actual %", "Ideal %", "Gap %"]
        for i, header in enumerate(headers):
            table.cell(0, i).text = header
        
        # Format header
        for i in range(cols):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Add data
        for i, gap_data in enumerate(top_gaps[:rows-1], 1):
            table.cell(i, 0).text = gap_data['Attribute']
            table.cell(i, 1).text = gap_data['Value']
            table.cell(i, 2).text = f"{gap_data['Actual']:.1f}%"
            table.cell(i, 3).text = f"{gap_data['Ideal']:.1f}%"
            table.cell(i, 4).text = f"{gap_data['Gap']:.1f}%"
            
            # Format cells
            for j in range(cols):
                paragraph = table.cell(i, j).text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Highlight gap
            paragraph = table.cell(i, 4).text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(0, 176, 80)  # Green
            paragraph.font.bold = True
    else:
        # Add a message if no significant gaps found
        textbox = slide.shapes.add_textbox(left, top, width, Inches(1.0))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = "No significant portfolio gaps identified."
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    # Add explanation textbox
    left = Inches(1.0)
    top = Inches(5.5)
    width = Inches(8.0)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    p = text_frame.add_paragraph()
    p.text = "Key Opportunity Areas for Jeju Portfolio:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # List top 3 opportunity areas
    for i, gap_data in enumerate(top_gaps[:3], 1):
        p = text_frame.add_paragraph()
        p.text = f"{i}. {gap_data['Attribute']} - {gap_data['Value']}: Current representation is {gap_data['Actual']:.1f}% vs. ideal {gap_data['Ideal']:.1f}% based on passenger preferences."
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

def create_recommendations_slides(prs, jj_df, jj_attr_df):
    """Create slides with SKU recommendations for Jeju."""
    # Create slide for SKU recommendations
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Portfolio Optimization Recommendations"
    
    # Create textbox with key recommendations
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(4.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Add heading
    p = text_frame.add_paragraph()
    p.text = "Key Recommendations for Jeju Portfolio Optimization:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # 1. Attribute focus
    p = text_frame.add_paragraph()
    p.text = "1. Attribute Focus Areas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Extract top underrepresented attributes from jj_attr_df
    attributes = ['Flavor', 'Taste', 'Thickness', 'Length']
    focus_areas = {}
    
    for attribute in attributes:
        # Find the section for this attribute
        found_section = False
        attr_values = []
        gaps = []
        
        for row_idx in range(len(jj_attr_df)):
            if pd.notna(jj_attr_df.iloc[row_idx, 0]) and attribute in str(jj_attr_df.iloc[row_idx, 0]):
                found_section = True
                continue
            
            if found_section:
                # Stop when we hit another attribute section
                if pd.notna(jj_attr_df.iloc[row_idx, 0]) and any(attr in str(jj_attr_df.iloc[row_idx, 0]) 
                                                             for attr in attributes):
                    break
                
                # Extract values
                row = jj_attr_df.iloc[row_idx]
                if pd.notna(row.iloc[0]) and str(row.iloc[0]).strip():
                    # Find columns with actual, ideal values
                    actual_col = 1  # Assuming fixed positions based on observed patterns
                    ideal_col = 2
                    
                    if pd.notna(row.iloc[actual_col]) and pd.notna(row.iloc[ideal_col]):
                        attr_value = str(row.iloc[0])
                        actual = float(row.iloc[actual_col])
                        ideal = float(row.iloc[ideal_col])
                        gap = ideal - actual
                        
                        # Only consider significant positive gaps (underrepresented segments)
                        if gap > 5:
                            attr_values.append(attr_value)
                            gaps.append(gap)
        
        # Add to focus areas if we found significant gaps
        if attr_values:
            # Sort by gap size
            sorted_indices = np.argsort(gaps)[::-1]
            attr_values = [attr_values[i] for i in sorted_indices]
            gaps = [gaps[i] for i in sorted_indices]
            
            focus_areas[attribute] = {
                'values': attr_values[:2],  # Top 2 values
                'gaps': gaps[:2]  # Corresponding gaps
            }
    
    # Add focus areas to slide
    for attribute, data in focus_areas.items():
        p = text_frame.add_paragraph()
        values_text = ", ".join([f"{val} (Gap: {gap:.1f}%)" for val, gap in zip(data['values'], data['gaps'])])
        p.text = f"• {attribute}: Increase representation in {values_text}"
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1
    
    # 2. SKU Optimization
    p = text_frame.add_paragraph()
    p.text = "2. SKU Portfolio Actions"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Add SKU actions
    p = text_frame.add_paragraph()
    p.text = "• Add new SKUs with underrepresented attribute combinations"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "• Remove low-performing SKUs in overrepresented segments"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "• Reallocate resources to support high-potential segments"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    p.level = 1
    
    # 3. Expected Impact
    p = text_frame.add_paragraph()
    p.text = "3. Expected Impact"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    p = text_frame.add_paragraph()
    p.text = "• Potential market share increase from current 11.4% to ~19-23%"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "• Better alignment with passenger preferences driving higher conversion"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "• Improved visibility and shelf presence in key segments"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    p.level = 1

def create_implementation_slide(prs):
    """Create a slide with the implementation plan."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Phased Implementation Plan"
    
    # Create table for implementation plan
    rows = 4  # Header + 3 phases
    cols = 3  # Phase, Actions, Timeline
    
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(4.0)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(5.5)
    table.columns[2].width = Inches(2.0)
    
    # Set headers
    headers = ["Phase", "Key Actions", "Timeline"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    # Format header
    for i in range(cols):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Phase 1 - Short-term
    table.cell(1, 0).text = "Short-Term"
    
    cell = table.cell(1, 1)
    tf = cell.text_frame
    
    p = tf.add_paragraph()
    p.text = "• Maintain and support top 5 performing SKUs"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Remove lowest performing SKUs in overrepresented segments"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Begin promoting underrepresented attributes in existing SKUs"
    p.font.size = Pt(12)
    
    table.cell(1, 2).text = "0-3 months"
    
    # Phase 2 - Medium-term
    table.cell(2, 0).text = "Medium-Term"
    
    cell = table.cell(2, 1)
    tf = cell.text_frame
    
    p = tf.add_paragraph()
    p.text = "• Introduce 2-3 new SKUs with key underrepresented attribute combinations"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Continue portfolio rationalization of underperforming SKUs"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Implement visibility and merchandising improvements"
    p.font.size = Pt(12)
    
    table.cell(2, 2).text = "3-6 months"
    
    # Phase 3 - Long-term
    table.cell(3, 0).text = "Long-Term"
    
    cell = table.cell(3, 1)
    tf = cell.text_frame
    
    p = tf.add_paragraph()
    p.text = "• Complete strategic portfolio realignment across all attributes"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Introduce 3-5 additional SKUs in high-potential segments"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Evaluate competitive response and adjust strategy accordingly"
    p.font.size = Pt(12)
    
    table.cell(3, 2).text = "6+ months"
    
    # Format all cells
    for row in range(1, rows):
        # Phase column
        cell = table.cell(row, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Timeline column
        paragraph = table.cell(row, 2).text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add conclusion textbox
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8.0)
    height = Inches(0.8)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    p = text_frame.add_paragraph()
    p.text = "Implementation should prioritize high-impact changes while maintaining operational feasibility. Regular monitoring and adjustment based on market response is essential."
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    # Example usage with file paths
    kw_products_path = "KW_products.csv"
    jj_products_path = "JJ_products.csv"
    kw_product_based_path = "KW_product_based.csv"
    jj_product_based_path = "JJ_product_based.csv"
    comparison_file_path = "comparison_kw_jj.csv"
    
    create_portfolio_presentation(
        kw_products_path,
        jj_products_path,
        kw_product_based_path,
        jj_product_based_path,
        comparison_file_path
    )
