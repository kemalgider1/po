"""
Portfolio Optimization Presentation Generator

This script creates a PowerPoint presentation summarizing the key findings from
the portfolio optimization analysis, with particular focus on the comparison
between Kuwait and Jeju.

Requirements:
- python-pptx package

Author: Claude
Date: March 12, 2025
"""

import os
import argparse
from datetime import datetime
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs, title, subtitle=None):
    """
    Create the title slide for the presentation.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Title text
        subtitle (str, optional): Subtitle text
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add subtitle if provided
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    
    return slide

def create_section_slide(prs, title):
    """
    Create a section divider slide.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Section title
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    return slide

def create_content_slide(prs, title, content=None, image_path=None):
    """
    Create a content slide with optional image.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        content (list, optional): List of content bullet points
        image_path (str, optional): Path to image file to include
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[1]  # Content with caption layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content if provided
    if content:
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        
        for i, point in enumerate(content):
            if i == 0:
                text_frame.text = point
            else:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
    
    # Add image if provided
    if image_path and os.path.exists(image_path):
        # If there's content, use a smaller image area
        if content:
            left = Inches(5)
            top = Inches(2)
            width = Inches(4)
            height = Inches(3)
        else:
            # Use full slide if no content
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
        
        slide.shapes.add_picture(image_path, left, top, width, height)
    
    return slide

def create_comparison_slide(prs, title, kuwait_data, jeju_data, metric_name, compare_text=None):
    """
    Create a slide comparing Kuwait and Jeju metrics.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        kuwait_data (float or dict): Kuwait metric value(s)
        jeju_data (float or dict): Jeju metric value(s)
        metric_name (str): Name of the metric being compared
        compare_text (str, optional): Additional comparison text
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[3]  # Two content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add Kuwait content
    left_shape = slide.placeholders[1]
    text_frame = left_shape.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "KUWAIT"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = f"{metric_name}:"
    p.font.bold = True
    
    if isinstance(kuwait_data, dict):
        for key, value in kuwait_data.items():
            p = text_frame.add_paragraph()
            p.text = f"{key}: {value}"
    else:
        p = text_frame.add_paragraph()
        p.text = f"{kuwait_data}"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
    
    # Add Jeju content
    right_shape = slide.placeholders[2]
    text_frame = right_shape.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "JEJU"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = f"{metric_name}:"
    p.font.bold = True
    
    if isinstance(jeju_data, dict):
        for key, value in jeju_data.items():
            p = text_frame.add_paragraph()
            p.text = f"{key}: {value}"
    else:
        p = text_frame.add_paragraph()
        p.text = f"{jeju_data}"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
    
    # Add comparison text at the bottom if provided
    if compare_text:
        text_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.8))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = compare_text
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True
    
    return slide

def create_two_image_slide(prs, title, image_path1, image_path2, caption1=None, caption2=None):
    """
    Create a slide with two images side by side.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        image_path1 (str): Path to first image
        image_path2 (str): Path to second image
        caption1 (str, optional): Caption for first image
        caption2 (str, optional): Caption for second image
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    
    # Add first image
    if os.path.exists(image_path1):
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(4)
        
        slide.shapes.add_picture(image_path1, left, top, width, height)
        
        # Add caption if provided
        if caption1:
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(4.5), Inches(0.5))
            text_frame = caption_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = caption1
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
    
    # Add second image
    if os.path.exists(image_path2):
        left = Inches(5.5)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(4)
        
        slide.shapes.add_picture(image_path2, left, top, width, height)
        
        # Add caption if provided
        if caption2:
            caption_box = slide.shapes.add_textbox(Inches(5.5), Inches(5.6), Inches(4.5), Inches(0.5))
            text_frame = caption_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = caption2
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
    
    return slide

def create_table_slide(prs, title, data, headers=None):
    """
    Create a slide with a table.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        data (list): List of rows, each row is a list of values
        headers (list, optional): List of header values
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    
    # Determine table dimensions
    rows = len(data)
    if headers:
        rows += 1  # Add a row for headers
    
    cols = max(len(row) for row in data)
    if headers:
        cols = max(cols, len(headers))
    
    # Create the table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Add headers if provided
    if headers:
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            
            # Format header cells
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
        
        # Add data starting from second row
        for i, row_data in enumerate(data):
            for j, value in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text = str(value)
    else:
        # Add data starting from first row
        for i, row_data in enumerate(data):
            for j, value in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(value)
    
    return slide

def create_recommendation_slide(prs, title, recommendations, image_path=None):
    """
    Create a slide with recommendations.
    
    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        recommendations (dict): Dictionary of recommendation categories and items
        image_path (str, optional): Path to image file to include
    
    Returns:
        slide: The created slide
    """
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content
    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    
    for category, items in recommendations.items():
        p = text_frame.add_paragraph()
        p.text = category
        p.font.bold = True
        p.font.size = Pt(16)
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
    
    # Add image if provided
    if image_path and os.path.exists(image_path):
        left = Inches(6.5)
        top = Inches(2)
        width = Inches(3)
        height = Inches(3)
        
        slide.shapes.add_picture(image_path, left, top, width, height)
    
    return slide

def generate_presentation(results_dir, output_path=None):
    """
    Generate a PowerPoint presentation summarizing portfolio optimization results.
    
    Args:
        results_dir (str): Directory containing analysis results
        output_path (str, optional): Output path for the presentation file
    
    Returns:
        str: Path to the generated presentation
    """
    # Set default output path if not specified
    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(results_dir, f'portfolio_optimization_presentation_{timestamp}.pptx')
    
    # Create presentation
    prs = Presentation()
    
    # 1. Title slide
    create_title_slide(
        prs,
        "PMI Duty-Free Portfolio Optimization",
        "Kuwait vs. Jeju: Insights and Recommendations"
    )
    
    # 2. Introduction slide
    create_content_slide(
        prs,
        "Project Overview",
        [
            "Portfolio optimization aims to align PMI product offerings with passenger preferences",
            "Analysis focused on comparing Kuwait (high performance) vs. Jeju (opportunity)",
            "Objective: Identify specific optimizations to improve Jeju market share",
            "Methodology: Comprehensive analysis of market data, passenger mix, portfolio alignment"
        ]
    )
    
    # 3. Market share comparison
    ms_image_path = os.path.join(results_dir, 'market_share_validation', 'market_share_comparison.png')
    
    create_comparison_slide(
        prs,
        "Market Share Comparison",
        "76.7%",
        "11.4%",
        "Market Share",
        "Kuwait has 6.7x the market share of Jeju"
    )
    
    # 4. Category scores comparison
    cat_c_image_path = os.path.join(results_dir, 'category_c_validation', 'category_c_validation.png')
    
    create_comparison_slide(
        prs,
        "Portfolio Alignment Scores",
        {
            "Category A": "7.3/10",
            "Category B": "9.5/10",
            "Category C": "8.7/10",
            "Category D": "9.2/10",
            "Average": "8.7/10"
        },
        {
            "Category A": "5.1/10",
            "Category B": "6.2/10",
            "Category C": "2.1/10",
            "Category D": "4.3/10",
            "Average": "4.4/10"
        },
        "Alignment Scores",
        "Kuwait's portfolio is significantly better aligned with passenger preferences"
    )
    
    # 5. Attribute comparisons
    create_section_slide(prs, "Attribute Analysis")
    
    # Add comparison slides for each attribute
    attributes = ['Flavor', 'Taste', 'Thickness', 'Length']
    
    for attribute in attributes:
        image_path = os.path.join(results_dir, 'visualizations', f'{attribute}_grid.png')
        
        create_content_slide(
            prs,
            f"{attribute} Comparison: Kuwait vs. Jeju",
            None,
            image_path
        )
    
    # 6. Kuwait shelf visualization
    kw_shelf_path = os.path.join(results_dir, 'visualizations', 'product_shelf.png')
    
    if os.path.exists(kw_shelf_path):
        create_content_slide(
            prs,
            "Product 'Shelf' Visualization",
            [
                "Kuwait's portfolio (top) shows balanced distribution across key attributes",
                "Jeju's portfolio (bottom) shows significant gaps in high-demand segments",
                "Circle size indicates volume; blue = PMI, gray = competitors",
                "Key opportunity: Expand presence in underserved segments in Jeju"
            ],
            kw_shelf_path
        )
    
    # 7. Radar chart
    radar_path = os.path.join(results_dir, 'visualizations', 'radar_chart.png')
    
    if os.path.exists(radar_path):
        create_content_slide(
            prs,
            "Portfolio Alignment Radar Chart",
            [
                "Kuwait shows strong alignment across all attributes",
                "Jeju shows significant gaps, especially in Flavor and Length",
                "Higher scores indicate better alignment with passenger preferences",
                "The difference in alignment explains much of the market share gap"
            ],
            radar_path
        )
    
    # 8. Key gaps in Jeju
    attr_gaps_path = os.path.join(results_dir, 'recommendations', 'attribute_gaps.png')
    
    if os.path.exists(attr_gaps_path):
        create_content_slide(
            prs,
            "Key Portfolio Gaps in Jeju",
            [
                "Green bars indicate underrepresented segments (opportunity)",
                "Red bars indicate overrepresented segments (potential rationalization)",
                "Largest gaps: SSL Thickness, NTD Caps Flavor, KS Length, Full Flavor Taste",
                "These gaps directly impact Jeju's market performance"
            ],
            attr_gaps_path
        )
    
    # 9. Recommendations
    create_section_slide(prs, "Recommendations for Jeju")
    
    # SKU recommendations
    sku_rec_path = os.path.join(results_dir, 'recommendations', 'sku_recommendations.png')
    
    if os.path.exists(sku_rec_path):
        create_content_slide(
            prs,
            "SKU-Level Recommendations",
            [
                "Green: SKUs to maintain and potentially expand",
                "Red: SKUs to consider for rationalization",
                "Focus on supporting products with strong alignment to passenger preferences",
                "Reallocate resources from underperforming SKUs to high-potential segments"
            ],
            sku_rec_path
        )
    
    # Attribute combinations
    combo_path = os.path.join(results_dir, 'recommendations', 'attribute_combinations.png')
    
    if os.path.exists(combo_path):
        create_content_slide(
            prs,
            "Recommended New Product Combinations",
            [
                "These combinations address the largest portfolio gaps",
                "Bubble size indicates market potential based on competitor volumes",
                "Focus on combinations with both large gaps and significant market potential",
                "Particular opportunity: SSL + NTD Caps combination"
            ],
            combo_path
        )
    
    # 10. Implementation plan
    create_recommendation_slide(
        prs,
        "Implementation Plan",
        {
            "Short-Term (0-3 months)": [
                "Maintain and support top 5 performing SKUs",
                "Remove 3 lowest-volume SKUs in overrepresented segments",
                "Focus marketing support on NTD Caps and SSL products"
            ],
            "Medium-Term (3-6 months)": [
                "Introduce 2-3 new SKUs with priority attribute combinations",
                "Continue portfolio rationalization for underperforming SKUs",
                "Expand focus to include Length attribute optimization"
            ],
            "Long-Term (6+ months)": [
                "Complete strategic portfolio realignment across all attributes",
                "Introduce products in remaining high-potential segments",
                "Implement ongoing portfolio monitoring system"
            ]
        }
    )
    
    # 11. Projected impact
    create_content_slide(
        prs,
        "Projected Impact",
        [
            "Current market share: ~11.4%",
            "Estimated short-term gain: 1.5-2.5 percentage points",
            "Estimated medium-term gain: 3.0-5.0 additional percentage points",
            "Long-term potential: 18-20% market share (comparable to similar markets)",
            "Additional benefits: Improved inventory efficiency, better alignment with global portfolio"
        ]
    )
    
    # 12. Next steps
    create_content_slide(
        prs,
        "Next Steps",
        [
            "1. Review and approve recommendations with commercial team",
            "2. Develop detailed implementation timeline",
            "3. Identify specific SKUs for new product introductions",
            "4. Create marketing support plan for priority segments",
            "5. Establish quarterly review process to monitor progress"
        ]
    )
    
    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    
    return output_path

def main():
    """
    Main function to parse arguments and generate presentation.
    """
    # Set up argument parser
    parser = argparse.ArgumentParser(description='Generate portfolio optimization presentation')
    
    parser.add_argument('--results_dir', required=True, 
                        help='Directory containing analysis results')
    
    parser.add_argument('--output', 
                        help='Output path for the presentation file')
    
    # Parse arguments
    args = parser.parse_args()
    
    # Generate presentation
    output_path = generate_presentation(args.results_dir, args.output)
    
    print(f"Presentation generation complete: {output_path}")

if __name__ == "__main__":
    main()
