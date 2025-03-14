"""
Portfolio Optimization Presentation Generator

This script creates a PowerPoint presentation based on the visualization
outputs from the portfolio visualization and recommendations engine scripts.
"""

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import pandas as pd
import numpy as np
import os
from pathlib import Path
import datetime


def create_title_slide(prs, title, subtitle):
    """
    Create the title slide for the presentation.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Presentation title
        subtitle (str): Presentation subtitle

    Returns:
        slide: The created slide
    """
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # Set title and subtitle
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Format title
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 0, 128)  # Navy Blue

    # Format subtitle
    subtitle_text_frame = subtitle_shape.text_frame
    subtitle_paragraph = subtitle_text_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(24)
    subtitle_paragraph.font.italic = True

    # Add date
    date_shape = slide.shapes.add_textbox(Inches(9), Inches(6.5), Inches(1), Inches(0.5))
    date_text_frame = date_shape.text_frame
    date_paragraph = date_text_frame.add_paragraph()
    date_paragraph.text = datetime.datetime.now().strftime("%B %d, %Y")
    date_paragraph.font.size = Pt(12)
    date_paragraph.font.italic = True

    return slide


def create_section_header(prs, title, subtitle=None):
    """
    Create a section header slide.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Section title
        subtitle (str, optional): Section subtitle

    Returns:
        slide: The created slide
    """
    # Add a section header slide
    section_layout = prs.slide_layouts[2]  # Section Header layout
    slide = prs.slides.add_slide(section_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Format title
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(40)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue

    # Add subtitle if provided
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

        # Format subtitle
        subtitle_text_frame = subtitle_shape.text_frame
        subtitle_paragraph = subtitle_text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(24)
        subtitle_paragraph.font.italic = True

    return slide


def create_content_slide(prs, title, content_paragraphs):
    """
    Create a content slide with bullet points.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        content_paragraphs (list): List of paragraph texts

    Returns:
        slide: The created slide
    """
    # Add a bullet slide
    bullet_slide_layout = prs.slide_layouts[1]  # Bullet layout
    slide = prs.slides.add_slide(bullet_slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Format title
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue

    # Add bullet points
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame

    # Clear any existing paragraphs
    if text_frame.paragraphs[0].text:
        text_frame.clear()

    # Add new paragraphs
    for i, text in enumerate(content_paragraphs):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.text = text
        paragraph.level = 0
        paragraph.font.size = Pt(20)

    return slide


def create_image_slide(prs, title, image_path, caption=None):
    """
    Create a slide with an image.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        image_path (str): Path to the image file
        caption (str, optional): Image caption

    Returns:
        slide: The created slide
    """
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Add image
    try:
        if os.path.exists(image_path):
            # Calculate image size to fit slide while maintaining aspect ratio
            img_width = Inches(9)
            top = Inches(1.5)
            left = Inches(0.5)

            # Add image
            slide.shapes.add_picture(image_path, left, top, width=img_width)

            # Add caption if provided
            if caption:
                caption_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
                caption_text_frame = caption_shape.text_frame
                caption_paragraph = caption_text_frame.add_paragraph()
                caption_paragraph.text = caption
                caption_paragraph.font.size = Pt(14)
                caption_paragraph.font.italic = True
                caption_paragraph.alignment = PP_ALIGN.CENTER
        else:
            # If image doesn't exist, add a placeholder
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(6), Inches(4))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)

            # Add text to placeholder
            text_frame = placeholder.text_frame
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"Image not found: {os.path.basename(image_path)}"
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER
    except Exception as e:
        # If there's an error loading the image, add an error message
        error_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        error_text_frame = error_shape.text_frame
        error_paragraph = error_text_frame.add_paragraph()
        error_paragraph.text = f"Error loading image: {str(e)}"
        error_paragraph.font.size = Pt(14)
        error_paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
        error_paragraph.alignment = PP_ALIGN.CENTER

    return slide


def create_comparison_slide(prs, title, left_title, right_title, left_image_path, right_image_path):
    """
    Create a slide comparing two images side by side.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        left_title (str): Title for left image
        right_title (str): Title for right image
        left_image_path (str): Path to left image
        right_image_path (str): Path to right image

    Returns:
        slide: The created slide
    """
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Add left image title
    left_title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.5))
    left_title_text_frame = left_title_shape.text_frame
    left_title_paragraph = left_title_text_frame.add_paragraph()
    left_title_paragraph.text = left_title
    left_title_paragraph.font.size = Pt(20)
    left_title_paragraph.font.bold = True
    left_title_paragraph.alignment = PP_ALIGN.CENTER

    # Add right image title
    right_title_shape = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(0.5))
    right_title_text_frame = right_title_shape.text_frame
    right_title_paragraph = right_title_text_frame.add_paragraph()
    right_title_paragraph.text = right_title
    right_title_paragraph.font.size = Pt(20)
    right_title_paragraph.font.bold = True
    right_title_paragraph.alignment = PP_ALIGN.CENTER

    # Add left image
    try:
        if os.path.exists(left_image_path):
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(4)
            slide.shapes.add_picture(left_image_path, left, top, width=width)
        else:
            # Placeholder for missing left image
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(4), Inches(3))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)

            text_frame = placeholder.text_frame
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"Image not found: {os.path.basename(left_image_path)}"
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER
    except Exception as e:
        # Error placeholder for left image
        error_shape = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4), Inches(1))
        error_text_frame = error_shape.text_frame
        error_paragraph = error_text_frame.add_paragraph()
        error_paragraph.text = f"Error loading image: {str(e)}"
        error_paragraph.font.size = Pt(14)
        error_paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
        error_paragraph.alignment = PP_ALIGN.CENTER

    # Add right image
    try:
        if os.path.exists(right_image_path):
            left = Inches(5.5)
            top = Inches(1.8)
            width = Inches(4)
            slide.shapes.add_picture(right_image_path, left, top, width=width)
        else:
            # Placeholder for missing right image
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.8), Inches(4), Inches(3))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)

            text_frame = placeholder.text_frame
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"Image not found: {os.path.basename(right_image_path)}"
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER
    except Exception as e:
        # Error placeholder for right image
        error_shape = slide.shapes.add_textbox(Inches(5.5), Inches(3), Inches(4), Inches(1))
        error_text_frame = error_shape.text_frame
        error_paragraph = error_text_frame.add_paragraph()
        error_paragraph.text = f"Error loading image: {str(e)}"
        error_paragraph.font.size = Pt(14)
        error_paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
        error_paragraph.alignment = PP_ALIGN.CENTER

    return slide


def create_table_slide(prs, title, data, columns=None):
    """
    Create a slide with a table of data.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        data (list or DataFrame): Table data
        columns (list, optional): Column headers

    Returns:
        slide: The created slide
    """
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Convert DataFrame to list if needed
    if isinstance(data, pd.DataFrame):
        if columns is None:
            columns = data.columns.tolist()
        data = [data.columns.tolist()] + data.values.tolist()
    elif columns is not None:
        data = [columns] + data

    # Calculate table dimensions
    rows = len(data)
    cols = len(data[0]) if rows > 0 else 0

    if rows > 0 and cols > 0:
        # Create table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        col_width = width / cols
        for i in range(cols):
            table.columns[i].width = col_width

        # Add data to table
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)

                # Format header row
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 76, 153)  # Blue
                    cell_paragraph = cell.text_frame.paragraphs[0]
                    cell_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                    cell_paragraph.font.bold = True

                # Format all cells
                cell_paragraph = cell.text_frame.paragraphs[0]
                cell_paragraph.font.size = Pt(12)
                cell_paragraph.alignment = PP_ALIGN.CENTER
    else:
        # Add a message if no data
        message_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        message_text_frame = message_shape.text_frame
        message_paragraph = message_text_frame.add_paragraph()
        message_paragraph.text = "No data available for table"
        message_paragraph.font.size = Pt(14)
        message_paragraph.font.italic = True
        message_paragraph.alignment = PP_ALIGN.CENTER

    return slide


def create_conclusion_slide(prs, title, recommendations):
    """
    Create a conclusion slide with key recommendations.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        recommendations (list): List of recommendation texts

    Returns:
        slide: The created slide
    """
    # Add a content slide layout
    content_slide_layout = prs.slide_layouts[1]  # Content layout
    slide = prs.slides.add_slide(content_slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Format title
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue

    # Add recommendations as bullet points
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame

    # Clear any existing paragraphs
    if text_frame.paragraphs[0].text:
        text_frame.clear()

    # Add recommendation paragraphs
    for i, rec in enumerate(recommendations):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.text = rec
        paragraph.level = 0
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(0, 102, 0)  # Dark Green

    return slide


def create_next_steps_slide(prs, title, next_steps):
    """
    Create a next steps slide.

    Args:
        prs (Presentation): PowerPoint presentation object
        title (str): Slide title
        next_steps (list): List of next step texts

    Returns:
        slide: The created slide
    """
    # Add a content slide layout
    content_slide_layout = prs.slide_layouts[1]  # Content layout
    slide = prs.slides.add_slide(content_slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Format title
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 76, 153)  # Blue

    # Add next steps as bullet points
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame

    # Clear any existing paragraphs
    if text_frame.paragraphs[0].text:
        text_frame.clear()

    # Add next step paragraphs
    for i, step in enumerate(next_steps):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.text = step
        paragraph.level = 0
        paragraph.font.size = Pt(24)

    return slide


def create_portfolio_optimization_presentation(visualization_dir, output_path):
    """
    Create a complete portfolio optimization presentation using visualization outputs.

    Args:
        visualization_dir (str): Directory containing visualization outputs
        output_path (str): Path to save the PowerPoint presentation

    Returns:
        str: Path to the created presentation
    """
    print(f"Creating portfolio optimization presentation...")

    # Create a blank presentation
    prs = Presentation()

    # 1. Create title slide
    create_title_slide(prs,
                       "Portfolio Optimization Project",
                       "Comparing Well-Aligned vs. Misaligned Markets: Kuwait and Jeju")

    # 2. Create introduction section
    create_section_header(prs, "Introduction", "Project Overview and Objectives")

    create_content_slide(prs, "Project Objectives", [
        "• Analyze PMI's product portfolio alignment with consumer preferences",
        "• Compare well-aligned markets (Kuwait) with misaligned markets (Jeju)",
        "• Identify optimization opportunities to improve market share",
        "• Generate specific SKU-level recommendations for portfolio adjustments",
        "• Develop an implementation plan for portfolio optimization"
    ])

    create_content_slide(prs, "Methodology", [
        "• Category C scoring to evaluate portfolio alignment with passenger mix",
        "• Attribute-based analysis across Flavor, Taste, Thickness, and Length",
        "• Comparison of actual vs. ideal product distribution",
        "• Identification of underrepresented and overrepresented segments",
        "• Development of visualizations for clear insights communication"
    ])

    # 3. Key Findings Section
    create_section_header(prs, "Key Findings", "Market Comparison Analysis")

    # Market share comparison
    create_image_slide(prs,
                       "Market Share Comparison",
                       os.path.join(visualization_dir, "market_share.png"),
                       "Kuwait (75%) vs. Jeju (12%) - Portfolio alignment directly impacts market share")

    # Portfolio alignment comparison
    create_image_slide(prs,
                       "Portfolio Alignment Radar Chart",
                       os.path.join(visualization_dir, "radar_chart.png"),
                       "Kuwait shows stronger alignment across all attributes compared to Jeju")

    # 4. Attribute Analysis Section
    create_section_header(prs, "Attribute Analysis", "Detailed Comparison by Product Attribute")

    # Side-by-side attribute comparisons
    create_comparison_slide(prs,
                            "Flavor Distribution: Kuwait vs. Jeju",
                            "Kuwait (Score: 9.64/10)",
                            "Jeju (Score: 7.53/10)",
                            os.path.join(visualization_dir, "Kuwait_Flavor_distribution.png"),
                            os.path.join(visualization_dir, "Jeju_Flavor_distribution.png"))

    create_comparison_slide(prs,
                            "Taste Distribution: Kuwait vs. Jeju",
                            "Kuwait (Score: 8.10/10)",
                            "Jeju (Score: 4.37/10)",
                            os.path.join(visualization_dir, "Kuwait_Taste_distribution.png"),
                            os.path.join(visualization_dir, "Jeju_Taste_distribution.png"))

    create_comparison_slide(prs,
                            "Thickness & Length Distribution",
                            "Kuwait (Aligned Portfolio)",
                            "Jeju (Misaligned Portfolio)",
                            os.path.join(visualization_dir, "Kuwait_Thickness_distribution.png"),
                            os.path.join(visualization_dir, "Jeju_Length_distribution.png"))

    # Add product shelf visualization
    create_image_slide(prs,
                       "Product 'Shelf' Visualization",
                       os.path.join(visualization_dir, "product_shelf.png"),
                       "Visual representation of current vs. ideal product distribution")

    # 5. Optimization Recommendations Section
    create_section_header(prs, "Optimization Recommendations", "SKU-Level Portfolio Adjustments")

    create_image_slide(prs,
                       "Jeju Optimization Recommendations",
                       os.path.join(visualization_dir, "Jeju_recommendations.png"),
                       "Specific SKU recommendations to address portfolio gaps")

    # Add top recommendations
    jeju_recommendations = [
        "1. Add new Menthol Caps SKUs to address 9.3% gap in Flavor attribute",
        "2. Add 1mg and Ultralights SKUs to address gaps in Taste attribute",
        "3. Adjust product mix to increase STD Thickness representation (80.6% gap)",
        "4. Remove duplicate Regular SKUs to reduce overrepresentation",
        "5. Prioritize Length attribute adjustments for quick portfolio alignment"
    ]

    create_content_slide(prs, "Top 5 Recommendations for Jeju", jeju_recommendations)

    # 6. Implementation Plan Section
    create_section_header(prs, "Implementation Plan", "Phased Approach to Portfolio Optimization")

    create_image_slide(prs,
                       "Implementation Plan for Jeju",
                       os.path.join(visualization_dir, "Jeju_implementation_plan.png"),
                       "Short, medium, and long-term actions for portfolio optimization")

    # Add implementation details
    implementation_details = [
        "Short-term (1-3 months):",
        "• Reduce overrepresentation by removing low-volume Regular SKUs",
        "• Adjust product attributes in existing SKUs where possible",
        "",
        "Medium-term (4-8 months):",
        "• Introduce new SKUs in underrepresented segments (Menthol, 1mg, STD)",
        "• Rebalance distribution focus across portfolio",
        "",
        "Long-term (8+ months):",
        "• Strategic portfolio review and alignment with passenger preferences",
        "• Development of optimization toolkit for all markets"
    ]

    create_content_slide(prs, "Phased Implementation Plan", implementation_details)

    # 7. Conclusion Section
    create_section_header(prs, "Conclusion", "Key Takeaways and Next Steps")

    # Add conclusion
    conclusions = [
        "✓ Portfolio alignment directly impacts market share performance",
        "✓ Kuwait demonstrates strong alignment across attributes (7.73/10)",
        "✓ Jeju shows significant misalignment in key attributes (6.02/10)",
        "✓ Specific SKU-level recommendations can address critical gaps",
        "✓ Phased implementation approach provides clear roadmap"
    ]

    create_conclusion_slide(prs, "Key Conclusions", conclusions)

    # Add next steps
    next_steps = [
        "1. Validate volume and margin impacts of recommendations",
        "2. Develop detailed implementation timeline with stakeholders",
        "3. Create simplified dashboard for ongoing portfolio monitoring",
        "4. Apply methodology to additional priority markets",
        "5. Integrate portfolio optimization with annual planning process"
    ]

    create_next_steps_slide(prs, "Next Steps", next_steps)

    # Save the presentation
    try:
        prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
        return output_path
    except Exception as e:
        print(f"Error saving presentation: {e}")
        return None


def main():
    """Main function to generate the portfolio optimization presentation"""
    # Define directories
    visualization_dir = "./visualization_results"
    optimization_dir = "./optimization_results"
    output_dir = "./presentations"

    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)

    # Define output path
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = os.path.join(output_dir, f"Portfolio_Optimization_Presentation_{timestamp}.pptx")

    # Create presentation
    create_portfolio_optimization_presentation(visualization_dir, output_path)

    print("Presentation generation completed!")


if __name__ == "__main__":
    main()
